#!/usr/bin/env python3
"""GovernorHub Document Sync — downloads targeted documents and rebuilds context.

Downloads only the folders that feed into the Ofsted Inspection Agent context,
then rebuilds combined_context.md with text extraction and compression.

Usage:
    python3 governorhub_sync.py              # Full sync + context rebuild
    python3 governorhub_sync.py --dry-run    # Show what would be downloaded
    python3 governorhub_sync.py --sync-only  # Download only, don't rebuild context
    python3 governorhub_sync.py --context-only  # Rebuild context from existing files
"""

from __future__ import annotations

import argparse
import json
import logging
import os
import re
import sys
import time
from datetime import datetime, timezone
from pathlib import Path

sys.path.insert(0, "/Users/timtrailor/code")
import credentials

# ── Config ──────────────────────────────────────────────────────────────────
BASE_URL = "https://app.governorhub.com"
DOCS_URL = f"{BASE_URL}/g/castlecefederation/docs"
LOGIN_URL = f"{BASE_URL}/login"
OUTPUT_DIR = Path.home() / "Desktop" / "School docs"
OFSTED_AGENT_DIR = Path.home() / "code" / "ofsted-agent"
CONTEXT_FILE = OFSTED_AGENT_DIR / "combined_context.md"
CONTEXT_FILE_ENC = OFSTED_AGENT_DIR / "combined_context.md.enc"
LOG_FILE = Path.home() / "code" / "governorhub_sync.log"
MAX_CONTEXT_TOKENS = 150_000

# Folders to sync — these feed the Ofsted agent context.
# Format: (folder_name, GovernorHub folder ID)
# IDs discovered via GraphQL introspection 2026-04-10.
TARGET_FOLDERS = [
    ("Ofsted 26 Victoria", "67fb7e1feff307c583e86f80"),
    ("Full Governing Body Meetings", "5ad097177b006d0a838b2356"),
    ("Admissions Committee Meetings", "5c3cd141d4ebd304a7c06c50"),
    ("Pupil & Curriculum Cttee Meetings", "5b3e098c65cd700006ff479c"),
    ("Resources Cttee Meetings", "5b3e0997eac6e900050fa7a8"),
    ("Governor Visits", "5af41dca310aee097a1b4400"),
    ("Safeguarding", None),  # IDs resolved at runtime
    ("SIAMS", None),
    ("Policies", "685b1d0856372f7690dc0961"),
    ("Helpful Documents", None),
    ("Risk Register", None),
    ("Training", None),
]

# Tier assignment for context building — order matters (Tier 1 first).
TIER_1_FOLDERS = {
    "KEY DOCS", "Governor Questions and Model Answers",
    "Linda (SEA) Ofsted Prep resources", "Victoria Assessment Details HFL 2025",
    "SEA Visit Reports",
}
TIER_1_ROOT = {"Ofsted 26 Victoria"}  # Root-level files in this folder = Tier 1

# Minutes/meetings: only include current + previous academic year in context
MINUTES_FOLDER = "Minutes"
MEETING_FOLDERS = {
    "Full Governing Body Meetings", "Admissions Committee Meetings",
    "Pupil & Curriculum Cttee Meetings", "Resources Cttee Meetings",
    "Governor Visits",
}
CURRENT_ACADEMIC_YEARS = ("2025-26", "2024-25")

# Boilerplate patterns to strip during text extraction
BOILERPLATE_RE = re.compile(
    r"(?:Table of Contents|CONTENTS|All rights reserved|"
    r"Copyright ©|Page \d+ of \d+|www\.\S+|https?://\S+)",
    re.IGNORECASE,
)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s %(levelname)s %(message)s",
    handlers=[
        logging.StreamHandler(),
        logging.FileHandler(LOG_FILE, encoding="utf-8"),
    ],
)
log = logging.getLogger("governorhub_sync")


# ═══════════════════════════════════════════════════════════════════════════
# PART 1: GovernorHub sync (Playwright + GraphQL)
# ═══════════════════════════════════════════════════════════════════════════

def sanitize_filename(name: str) -> str:
    return re.sub(r'[<>:"/\\|?*]', "_", name).strip()


def gql_query(page, query: str) -> dict:
    result = page.evaluate(
        """async (payload) => {
            const resp = await fetch('/api/graphql', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify(payload)
            });
            return await resp.json();
        }""",
        {"query": query},
    )
    if result.get("errors"):
        log.warning("GraphQL errors: %s", json.dumps(result["errors"])[:300])
    return result


def login(page):
    log.info("Logging in to GovernorHub...")
    page.goto(LOGIN_URL, wait_until="domcontentloaded")
    page.wait_for_timeout(3000)
    page.fill('input[type="email"]', credentials.GOVERNORHUB_EMAIL)
    page.click('button[type="button"]')
    page.wait_for_timeout(3000)
    page.fill('input[type="password"]', credentials.GOVERNORHUB_PASSWORD)
    page.click('button[type="submit"]')
    page.wait_for_timeout(5000)
    if "interaction" in page.url:
        raise RuntimeError(f"Login failed — still on {page.url}")
    log.info("Logged in successfully")


def resolve_folder_ids(page) -> dict:
    """Resolve folder IDs for TARGET_FOLDERS that don't have hardcoded IDs."""
    page.goto(DOCS_URL, wait_until="domcontentloaded")
    page.wait_for_timeout(5000)

    id_map = {}
    links = page.query_selector_all('tr a[href*="/docs/"]')
    for link in links:
        href = link.get_attribute("href") or ""
        name = link.inner_text().strip()
        m = re.search(r"/docs/([a-f0-9]{24})$", href)
        if m:
            id_map[name] = m.group(1)
    return id_map


def list_folder_contents(page, folder_id: str) -> list[dict]:
    query = """{
        file(id: "%s") {
            parentOfFilesConnection(first: 500) {
                edges { node { _id filename folder lastModified preventDownload } }
            }
        }
    }""" % folder_id
    result = gql_query(page, query)
    items = []
    edges = (result.get("data", {}).get("file", {})
             .get("parentOfFilesConnection", {}).get("edges", []))
    for edge in edges:
        n = edge["node"]
        items.append({
            "id": n["_id"], "filename": n["filename"],
            "is_folder": bool(n.get("folder")),
            "last_modified": n.get("lastModified"),
            "prevent_download": bool(n.get("preventDownload")),
        })
    return items


def get_download_url(page, file_id: str) -> str | None:
    query = '{ file(id: "%s") { fileViewer { signedUrl } } }' % file_id
    result = gql_query(page, query)
    return (result.get("data", {}).get("file", {})
            .get("fileViewer", {}).get("signedUrl"))


def download_file(page, file_id: str, filename: str, dest_dir: Path,
                  dry_run: bool, remote_modified: str = None) -> str:
    """Returns 'downloaded', 'updated', or 'skipped'."""
    dest_path = dest_dir / sanitize_filename(filename)

    if dest_path.exists() and remote_modified:
        try:
            remote_dt = datetime.fromisoformat(remote_modified.replace("Z", "+00:00"))
            local_mtime = datetime.fromtimestamp(dest_path.stat().st_mtime, tz=timezone.utc)
            if remote_dt <= local_mtime:
                return "skipped"
            if dry_run:
                log.info("  [DRY RUN] Would UPDATE: %s", filename)
                return "updated"
        except (ValueError, OSError):
            return "skipped"
    elif dest_path.exists():
        return "skipped"

    if dry_run:
        log.info("  [DRY RUN] Would download: %s", filename)
        return "downloaded"

    signed_url = get_download_url(page, file_id)
    if not signed_url:
        log.warning("  No download URL for: %s", filename)
        return "skipped"

    try:
        resp = page.request.get(signed_url, fail_on_status_code=False)
        if resp.status != 200:
            log.warning("  Download failed (HTTP %d): %s", resp.status, filename)
            return "skipped"
        body = resp.body()
        dest_dir.mkdir(parents=True, exist_ok=True)
        was_update = dest_path.exists()
        dest_path.write_bytes(body)
        if remote_modified:
            try:
                ts = datetime.fromisoformat(remote_modified.replace("Z", "+00:00")).timestamp()
                os.utime(dest_path, (ts, ts))
            except (ValueError, OSError):
                pass
        action = "Updated" if was_update else "Downloaded"
        log.info("  %s: %s (%s bytes)", action, filename, f"{len(body):,}")
        return "updated" if was_update else "downloaded"
    except Exception as e:
        log.error("  Error downloading %s: %s", filename, e)
        return "skipped"


def sync_folder_recursive(page, folder_id: str, folder_name: str, dest_dir: Path,
                          dry_run: bool, stats: dict, depth: int = 0,
                          year_filter: bool = False):
    indent = "  " * depth
    log.info("%sSyncing: %s", indent, folder_name)
    items = list_folder_contents(page, folder_id)
    for item in items:
        if item["is_folder"]:
            # For meeting folders, skip academic years older than current + previous
            if year_filter and depth == 0:
                if not any(y in item["filename"] for y in CURRENT_ACADEMIC_YEARS):
                    log.debug("%s  Skipping old year: %s", indent, item["filename"])
                    continue
            sync_folder_recursive(
                page, item["id"], item["filename"],
                dest_dir / sanitize_filename(item["filename"]),
                dry_run, stats, depth + 1,
                year_filter=year_filter,
            )
        else:
            if item["prevent_download"]:
                stats["skipped"] += 1
                continue
            stats["total"] += 1
            result = download_file(
                page, item["id"], item["filename"], dest_dir,
                dry_run, item.get("last_modified"),
            )
            stats[result] = stats.get(result, 0) + 1
    time.sleep(0.3)


def run_sync(dry_run: bool) -> dict:
    from playwright.sync_api import sync_playwright

    stats = {"total": 0, "downloaded": 0, "updated": 0, "skipped": 0}

    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page()
        page.set_default_timeout(30000)

        try:
            login(page)
            id_map = resolve_folder_ids(page)

            for folder_name, folder_id in TARGET_FOLDERS:
                fid = folder_id or id_map.get(folder_name)
                if not fid:
                    log.warning("Could not resolve ID for folder: %s", folder_name)
                    continue
                dest = OUTPUT_DIR / sanitize_filename(folder_name)
                is_meeting = folder_name in MEETING_FOLDERS
                sync_folder_recursive(page, fid, folder_name, dest, dry_run, stats,
                                      year_filter=is_meeting)
        finally:
            browser.close()

    log.info("Sync: %d new, %d updated, %d skipped (of %d total)",
             stats.get("downloaded", 0), stats.get("updated", 0),
             stats["skipped"], stats["total"])
    return stats


# ═══════════════════════════════════════════════════════════════════════════
# PART 2: Context builder — extract text, compress, assemble
# ═══════════════════════════════════════════════════════════════════════════

def extract_text(filepath: Path) -> str:
    """Extract plain text from a document file."""
    ext = filepath.suffix.lower()
    try:
        if ext == ".docx":
            return _extract_docx(filepath)
        elif ext == ".pdf":
            return _extract_pdf(filepath)
        elif ext == ".pptx":
            return _extract_pptx(filepath)
        elif ext in (".xlsx", ".xls"):
            return _extract_xlsx(filepath)
        elif ext in (".txt", ".md", ".csv", ".html", ".htm"):
            return filepath.read_text(errors="replace")
        elif ext in (".doc", ".odt"):
            return _extract_with_textutil(filepath)
        else:
            log.debug("  Skipping unsupported format: %s", filepath.name)
            return ""
    except Exception as e:
        log.warning("  Failed to extract %s: %s", filepath.name, e)
        return ""


def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_pdf(path: Path) -> str:
    import pdfplumber
    parts = []
    with pdfplumber.open(str(path)) as pdf:
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                parts.append(text)
    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), data_only=True)
    parts = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _extract_with_textutil(path: Path) -> str:
    """Use macOS textutil to convert .doc/.odt to text."""
    import subprocess
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as tmp:
        tmp_path = tmp.name
    try:
        subprocess.run(
            ["textutil", "-convert", "txt", "-output", tmp_path, str(path)],
            capture_output=True, timeout=30,
        )
        return Path(tmp_path).read_text(errors="replace")
    except Exception:
        return ""
    finally:
        Path(tmp_path).unlink(missing_ok=True)


def compress_text(text: str) -> str:
    """Remove boilerplate, collapse whitespace, strip noise."""
    # Remove boilerplate lines
    lines = text.split("\n")
    cleaned = []
    for line in lines:
        stripped = line.strip()
        if not stripped:
            # Collapse multiple blank lines to one
            if cleaned and cleaned[-1] != "":
                cleaned.append("")
            continue
        if BOILERPLATE_RE.search(stripped):
            continue
        # Skip very short lines that are likely headers/footers/page numbers
        if len(stripped) < 4 and not stripped[0].isalpha():
            continue
        cleaned.append(stripped)

    result = "\n".join(cleaned)
    # Collapse runs of 3+ newlines to 2
    result = re.sub(r"\n{3,}", "\n\n", result)
    return result.strip()


def estimate_tokens(text: str) -> int:
    """Rough token estimate: ~4 chars per token for English text."""
    return len(text) // 4


def assign_tier(folder_path: str, filename: str) -> int:
    """Assign a tier (1 or 2) based on the folder path."""
    parts = Path(folder_path).parts
    # Files directly in Ofsted 26 Victoria root or its Tier 1 subfolders
    if len(parts) >= 1 and parts[0] == "Ofsted 26 Victoria":
        if len(parts) == 1:
            return 1  # Root-level files
        if parts[1] in TIER_1_FOLDERS:
            return 1
    # Safeguarding current year
    if "Safeguarding" in parts and "2025-26" in folder_path:
        return 1
    return 2


def _classify_document(rel_folder: str, filename: str) -> tuple:
    """Classify a document into (section, meeting_date_or_none).

    Sections: 'ofsted_prep', 'fgb', 'resources', 'p_and_c', 'admissions',
              'safeguarding', 'siams', 'policies', 'other'
    """
    parts = Path(rel_folder).parts

    # Meeting folders — extract meeting date from folder name
    meeting_date = None
    for p in parts:
        # Match date-like folder names: "2026-03-25", "2025 03 26", etc.
        m = re.match(r"(\d{4})[-_ ](\d{2})[-_ ](\d{2})", p)
        if m:
            meeting_date = f"{m.group(1)}-{m.group(2)}-{m.group(3)}"

    # Committee meetings — check FIRST (even inside Ofsted 26 Victoria)
    if "Full Governing Body Meetings" in rel_folder:
        return ("fgb", meeting_date)
    if "Resources Cttee Meetings" in rel_folder:
        return ("resources", meeting_date)
    if "Pupil & Curriculum Cttee Meetings" in rel_folder:
        return ("p_and_c", meeting_date)
    if "Admissions Committee Meetings" in rel_folder:
        return ("admissions", meeting_date)
    if "Governor Visits" in rel_folder:
        return ("governor_visits", meeting_date)
    # Minutes inside Ofsted 26 Victoria — classify by committee
    if MINUTES_FOLDER in rel_folder:
        for keyword, section in [("FGB", "fgb"), ("Resources", "resources"),
                                 ("P&C", "p_and_c"), ("Admissions", "admissions")]:
            if keyword in rel_folder:
                return (section, meeting_date)
    # Non-meeting sections
    if "Ofsted 26 Victoria" in rel_folder:
        return ("ofsted_prep", None)
    if "Safeguarding" in rel_folder:
        return ("safeguarding", None)
    if "SIAMS" in rel_folder:
        return ("siams", None)
    if "Policies" in rel_folder:
        return ("policies", None)
    if "Helpful Documents" in rel_folder:
        return ("helpful", None)
    if "Risk Register" in rel_folder:
        return ("risk_register", None)
    if "Training" in rel_folder:
        return ("training", None)
    return ("other", None)


SECTION_LABELS = {
    "ofsted_prep": "OFSTED PREPARATION — Key Documents for Victoria Inspection",
    "fgb": "FULL GOVERNING BODY (FGB) MEETINGS",
    "resources": "RESOURCES COMMITTEE MEETINGS",
    "p_and_c": "PUPIL & CURRICULUM COMMITTEE MEETINGS",
    "admissions": "ADMISSIONS COMMITTEE MEETINGS",
    "governor_visits": "GOVERNOR VISITS",
    "safeguarding": "SAFEGUARDING",
    "siams": "SIAMS (Church Inspection)",
    "policies": "CURRENT POLICIES",
    "helpful": "HELPFUL DOCUMENTS",
    "risk_register": "RISK REGISTER",
    "training": "TRAINING",
    "other": "OTHER DOCUMENTS",
}

# Section ordering — most important first
SECTION_ORDER = [
    "ofsted_prep", "fgb", "resources", "p_and_c",
    "admissions", "governor_visits", "safeguarding", "siams",
    "policies", "helpful", "risk_register", "training", "other",
]

# Committee sections that get grouped by meeting date
COMMITTEE_SECTIONS = {"fgb", "resources", "p_and_c", "admissions", "governor_visits"}


def build_context() -> str:
    """Walk downloaded folders, extract text, build combined_context.md with clear sections."""
    log.info("Building context from downloaded files...")

    # Collect all documents, classified by section
    sections = {s: [] for s in SECTION_ORDER}

    for folder_name, _ in TARGET_FOLDERS:
        folder_path = OUTPUT_DIR / sanitize_filename(folder_name)
        if not folder_path.exists():
            log.warning("Folder not found: %s", folder_path)
            continue

        for filepath in sorted(folder_path.rglob("*")):
            if filepath.is_dir():
                continue

            rel_path = filepath.relative_to(OUTPUT_DIR)
            rel_folder = str(rel_path.parent)

            # Minutes/meetings filter: only current + previous academic year
            is_meeting = MINUTES_FOLDER in rel_folder or any(
                mf in rel_folder for mf in MEETING_FOLDERS
            )
            if is_meeting:
                if not any(y in rel_folder for y in CURRENT_ACADEMIC_YEARS):
                    continue

            text = extract_text(filepath)
            if not text or len(text.strip()) < 50:
                continue

            text = compress_text(text)
            mtime = datetime.fromtimestamp(filepath.stat().st_mtime)
            section, meeting_date = _classify_document(rel_folder, filepath.name)

            entry = {
                "filename": filepath.name,
                "folder": rel_folder,
                "modified": mtime.strftime("%Y-%m-%d"),
                "meeting_date": meeting_date,
                "text": text,
                "tokens": estimate_tokens(text),
            }

            if section in sections:
                sections[section].append(entry)
            else:
                sections["other"].append(entry)

    # Sort within each section: meetings by meeting date (newest first),
    # others by modification date (newest first)
    for section, docs in sections.items():
        docs.sort(
            key=lambda d: d.get("meeting_date") or d["modified"],
            reverse=True,
        )

    # Build meeting date index for committee sections
    meeting_index_lines = []
    for section in COMMITTEE_SECTIONS:
        dates = sorted(set(
            d["meeting_date"] for d in sections[section] if d.get("meeting_date")
        ), reverse=True)
        if dates:
            label = SECTION_LABELS[section].split(" — ")[0]
            meeting_index_lines.append(
                f"- **{label}**: {', '.join(dates)}"
            )

    meeting_index = ""
    if meeting_index_lines:
        meeting_index = (
            "\n## MEETING DATES INDEX\n"
            "The following committee meetings have documents in this context "
            "(most recent first):\n"
            + "\n".join(meeting_index_lines)
            + "\n\nWhen asked about the 'last' or 'most recent' meeting, "
            "refer to this index to identify the correct date.\n"
        )

    # Assemble the context file
    header = f"""# Castle CE Federation — Ofsted Inspection Reference
# Victoria CE Infant & Nursery School
# Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}

This document contains key school documents for the Ofsted inspection.
Documents are organised by section. Each document is tagged with
[SOURCE: filename | folder | date] for reference.

Key schools:
- **Victoria CE Infant & Nursery School** — Last inspected Oct 2023 (Requires Improvement)
- **Thomas Coram CE School** — Last inspected June 2023 (Good)
- **Castle CE Federation** — Federation of both schools, shared governing body
{meeting_index}
"""

    # Allocate token budget per section — ensures every section gets represented.
    # Priority sections get more budget; remaining is split among the rest.
    SECTION_BUDGETS = {
        "ofsted_prep": 40_000,
        "fgb": 25_000,
        "resources": 15_000,
        "p_and_c": 15_000,
        "admissions": 5_000,
        "governor_visits": 10_000,
        "safeguarding": 15_000,
        "siams": 5_000,
        "policies": 10_000,
        "helpful": 3_000,
        "risk_register": 3_000,
        "training": 3_000,
        "other": 1_000,
    }

    parts = [header]
    token_count = estimate_tokens(header)
    total_docs = 0

    for section in SECTION_ORDER:
        docs = sections[section]
        if not docs:
            continue

        label = SECTION_LABELS[section]
        section_budget = SECTION_BUDGETS.get(section, 5_000)
        section_used = 0
        section_header = f"\n{'=' * 70}\n# {label}\n{'=' * 70}\n"

        # For committee meetings, group by meeting date
        if section in COMMITTEE_SECTIONS:
            parts.append(section_header)
            token_count += 20
            section_used += 20

            by_date = {}
            undated = []
            for doc in docs:
                md = doc.get("meeting_date")
                if md:
                    by_date.setdefault(md, []).append(doc)
                else:
                    undated.append(doc)

            for meeting_date in sorted(by_date.keys(), reverse=True):
                if section_used >= section_budget:
                    break
                meeting_docs = by_date[meeting_date]
                meeting_header = f"\n--- Meeting: {meeting_date} ---\n"
                parts.append(meeting_header)
                token_count += 5
                section_used += 5

                for doc in meeting_docs:
                    doc_tokens = doc["tokens"] + 10
                    if section_used + doc_tokens > section_budget:
                        break
                    if token_count + doc_tokens > MAX_CONTEXT_TOKENS:
                        break
                    source_tag = (f"[SOURCE: {doc['filename']} | "
                                  f"{doc['folder']} | Modified: {doc['modified']}]")
                    parts.append(f"\n{source_tag}\n\n{doc['text']}\n")
                    token_count += doc_tokens
                    section_used += doc_tokens
                    total_docs += 1

            for doc in undated:
                if section_used >= section_budget:
                    break
                doc_tokens = doc["tokens"] + 10
                if token_count + doc_tokens > MAX_CONTEXT_TOKENS:
                    break
                source_tag = (f"[SOURCE: {doc['filename']} | "
                              f"{doc['folder']} | Modified: {doc['modified']}]")
                parts.append(f"\n---\n{source_tag}\n\n{doc['text']}\n")
                token_count += doc_tokens
                section_used += doc_tokens
                total_docs += 1
        else:
            parts.append(section_header)
            token_count += 20
            section_used += 20

            for doc in docs:
                doc_tokens = doc["tokens"] + 10
                if section_used + doc_tokens > section_budget:
                    break
                if token_count + doc_tokens > MAX_CONTEXT_TOKENS:
                    break
                source_tag = (f"[SOURCE: {doc['filename']} | "
                              f"{doc['folder']} | Modified: {doc['modified']}]")
                parts.append(f"\n---\n{source_tag}\n\n{doc['text']}\n")
                token_count += doc_tokens
                section_used += doc_tokens
                total_docs += 1

        log.info("  Section %s: %d tokens used (budget %d)",
                 section, section_used, section_budget)

    context = "".join(parts)

    actual_estimate = estimate_tokens(context)
    context = context.replace(
        "# Generated:",
        f"# Total tokens: ~{actual_estimate:,}\n# Generated:",
        1,
    )

    log.info("Context built: %d documents, ~%s tokens, %s chars",
             total_docs, f"{actual_estimate:,}", f"{len(context):,}")

    return context


def encrypt_context(plaintext: str):
    """Encrypt context with Fernet for the public GitHub repo."""
    from cryptography.fernet import Fernet
    key = credentials.CONTEXT_KEY.encode()
    encrypted = Fernet(key).encrypt(plaintext.encode())
    CONTEXT_FILE_ENC.write_bytes(encrypted)
    log.info("Encrypted context written to %s", CONTEXT_FILE_ENC)


def restart_streamlit():
    """Restart the governors Streamlit app so it picks up the new context."""
    import subprocess
    try:
        uid = subprocess.check_output(["id", "-u"], text=True).strip()
        subprocess.run(
            ["launchctl", "kickstart", "-k", f"gui/{uid}/com.timtrailor.governors"],
            capture_output=True, timeout=10,
        )
        log.info("Streamlit app restarted to load new context")
    except Exception as e:
        log.warning("Could not restart Streamlit: %s", e)


def rebuild_context():
    """Build, write, encrypt the context file, and restart Streamlit."""
    context = build_context()
    CONTEXT_FILE.write_text(context)
    log.info("Context written to %s (%s chars)", CONTEXT_FILE, f"{len(context):,}")
    encrypt_context(context)
    restart_streamlit()


# ═══════════════════════════════════════════════════════════════════════════
# PART 3: Main
# ═══════════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(description="GovernorHub sync + context builder")
    parser.add_argument("--dry-run", action="store_true",
                        help="Show what would be downloaded, don't rebuild context")
    parser.add_argument("--sync-only", action="store_true",
                        help="Download files only, don't rebuild context")
    parser.add_argument("--context-only", action="store_true",
                        help="Rebuild context from existing files, don't sync")
    args = parser.parse_args()

    log.info("=" * 60)
    log.info("GovernorHub Sync — %s", datetime.now().strftime("%Y-%m-%d %H:%M"))
    log.info("=" * 60)

    changed = False

    if not args.context_only:
        stats = run_sync(args.dry_run)
        changed = stats.get("downloaded", 0) + stats.get("updated", 0) > 0

    if not args.dry_run and not args.sync_only:
        if args.context_only or changed:
            rebuild_context()
        else:
            log.info("No files changed — skipping context rebuild")

    log.info("Done.")


if __name__ == "__main__":
    main()
